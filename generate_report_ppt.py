# generate_report.py
import os, json, datetime, base64
import pandas as pd
import plotly.express as px
from jinja2 import Environment, select_autoescape
from ai_utils import generate_text   # Gemini direct call
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
# ---------- SETTINGS ----------
DATA_FILE = "audit_data.json"
OUTPUT_HTML = "report.html"
COMPANY_NAME = "Chola MS Risk Services"
COMPANY_LOGO = "Chola_Risk.png"  # Logo file path
AUDIT_DATE = None
# ------------------------------

# ---------- HELPERS ----------
def load_data(path):
    with open(path, "r", encoding="utf-8") as f:
        return pd.DataFrame(json.load(f))

def encode_image_to_base64(image_path):
    """Convert image to base64 for embedding in HTML"""
    try:
        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            return f"data:image/png;base64,{encoded_string}"
    except FileNotFoundError:
        print(f"⚠️  Logo file not found: {image_path}")
        return None

def safe_get_col(df, name, default=""):
    return df[name] if name in df.columns else pd.Series([default]*len(df))

def compute_kpis(df):
    high = int((df["Priority"].str.upper() == "HIGH").sum()) if "Priority" in df else 0
    med = int((df["Priority"].str.upper() == "MEDIUM").sum()) if "Priority" in df else 0
    low = int((df["Priority"].str.upper() == "LOW").sum()) if "Priority" in df else 0
    return {"high": high, "med": med, "low": low, "total": len(df)}

def prepare_records(df):
    df["Sr. No"] = df.get("Sr. No", pd.Series(range(1, len(df)+1)))
    df["Observation"] = safe_get_col(df, "Observation", "")
    df["Recommendation"] = safe_get_col(df, "Recommendation", "")
    df["Priority"] = safe_get_col(df, "Priority", "")
    df["Priority_Score"] = pd.to_numeric(df.get("Priority_Score", 0), errors="coerce").fillna(0).astype(int)
    df["Type of Hazard"] = safe_get_col(df, "Type of Hazard", "")
    df["Location_Norm"] = safe_get_col(df, "Location_Norm", "")
    df["Risk_Score"] = pd.to_numeric(df.get("Risk_Score", 0), errors="coerce").fillna(0.0)
    df["Asset Category"] = safe_get_col(df, "Asset Category", "")
    return df

# ---------- AI SUMMARIES ----------
def generate_chart_summary(title, df):
    prompt = f"""
    You are an audit data analyst.
    Based on the chart '{title}' and the following dataset sample:

    {df.head(20).to_string()}

    Write 3-5 bullet point insights focusing on trends, risks, anomalies, and key percentages.
    """
    return generate_text(prompt, max_output_tokens=400)

def generate_exec_summary(metrics, examples):
    prompt = f"""
    You are an expert industrial safety auditor.
    Write a concise 3-paragraph executive summary for company leadership:
    Total findings: {metrics['total']}
    High: {metrics['high']}, Medium: {metrics['med']}, Low: {metrics['low']}
    Example top risks:
    """
    for ex in examples:
        prompt += f"- {ex.get('Sr. No')} | {ex.get('Risk_Score',0)} | {ex.get('Location_Norm','')} | {ex.get('Type of Hazard','')} | {ex.get('Observation','')[:80]}\n"
    prompt += "\nKeep it factual, under 120 words/para. End with 3 bullet-point action items."
    return generate_text(prompt, max_output_tokens=800)

def generate_priority_details(df, priority):
    if priority.upper() == "ALL":
        sub_df = df
    else:
        sub_df = df[df["Priority"].str.upper() == priority.upper()]
    prompt = f"""
    You are an expert auditor.
    Generate 3-5 key observations and 3-5 recommendations from the following dataset with key percentages:

    {sub_df.head(20).to_string()}
    """
    return generate_text(prompt, max_output_tokens=400)
def create_ppt_report(company_name, audit_date, kpis, exec_summary, priority_details, df, logo_path=None):
    """Generate PowerPoint presentation for audit report"""
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    primary_blue = RGBColor(30, 64, 175)
    danger_red = RGBColor(220, 38, 38)
    warning_orange = RGBColor(234, 88, 12)
    success_green = RGBColor(5, 150, 105)
    text_dark = RGBColor(31, 41, 55)
    
    # Slide 1: Title Slide
    create_title_slide(prs, company_name, audit_date, logo_path, primary_blue)
    
    # Slide 2: Executive Summary
    create_executive_summary_slide(prs, exec_summary, primary_blue, text_dark)
    
    # Slide 3: KPI Overview
    create_kpi_slide(prs, kpis, primary_blue, danger_red, warning_orange, success_green)
    
    # Slide 4: Risk Distribution
    create_risk_distribution_slide(prs, df, primary_blue, text_dark)
    
    # Slide 5: Priority Analysis - High
    create_priority_slide(prs, "HIGH PRIORITY RISKS", priority_details.get("HIGH", ""), 
                         danger_red, text_dark, df[df["Priority"].str.upper() == "HIGH"] if "Priority" in df.columns else df.head(0))
    
    # Slide 6: Priority Analysis - Medium
    create_priority_slide(prs, "MEDIUM PRIORITY RISKS", priority_details.get("MEDIUM", ""), 
                         warning_orange, text_dark, df[df["Priority"].str.upper() == "MEDIUM"] if "Priority" in df.columns else df.head(0))
    
    # Slide 7: Location Analysis
    create_location_analysis_slide(prs, df, primary_blue, text_dark)
    
    # Slide 8: Hazard Type Analysis
    create_hazard_analysis_slide(prs, df, primary_blue, text_dark)
    
    # Slide 9: Top Risk Findings
    create_top_findings_slide(prs, df, primary_blue, danger_red, text_dark)
    
    # Slide 10: Recommendations & Next Steps
    create_recommendations_slide(prs, df, kpis, primary_blue, text_dark)
    
    # Save presentation
    output_path = f"Chola_MS_Risk_Services_Audit_Report_{audit_date.replace('-', '_')}.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation generated: {output_path}")
    return output_path

def create_title_slide(prs, company_name, audit_date, logo_path, primary_blue):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add logo if available
    if logo_path and os.path.exists(logo_path):
        try:
            logo = slide.shapes.add_picture(logo_path, Inches(1), Inches(0.5), height=Inches(1.5))
        except:
            pass
    
    # Company name
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = company_name
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = primary_blue
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Comprehensive Risk Audit Report"
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = RGBColor(55, 65, 81)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(0.5))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = f"Generated on {audit_date}"
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = RGBColor(107, 114, 128)
    date_p.alignment = PP_ALIGN.CENTER

def create_executive_summary_slide(prs, exec_summary, primary_blue, text_dark):
    """Create executive summary slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    content.text = clean_html_text(exec_summary)
    
    # Format content text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_dark
        paragraph.space_after = Pt(12)

def create_kpi_slide(prs, kpis, primary_blue, danger_red, warning_orange, success_green):
    """Create KPI overview slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Key Performance Indicators"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = primary_blue
    title_p.alignment = PP_ALIGN.CENTER
    
    # KPI boxes
    kpi_data = [
        ("Critical Risks", kpis["high"], "Immediate Action", danger_red),
        ("Medium Risks", kpis["med"], "Monitor & Plan", warning_orange),
        ("Low Risks", kpis["low"], "Routine Maintenance", success_green),
        ("Total Findings", kpis["total"], "Complete Assessment", primary_blue)
    ]
    
    box_width = Inches(2.5)
    box_height = Inches(2)
    start_x = Inches(1)
    y_pos = Inches(2.5)
    spacing = Inches(0.3)
    
    for i, (label, value, desc, color) in enumerate(kpi_data):
        x_pos = start_x + i * (box_width + spacing)
        
        # Create rounded rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(249, 250, 251)
        shape.line.color.rgb = color
        shape.line.width = Pt(3)
        
        # Add text
        text_box = slide.shapes.add_textbox(x_pos, y_pos, box_width, box_height)
        text_frame = text_box.text_frame
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        
        # Label
        p1 = text_frame.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(75, 85, 99)
        p1.alignment = PP_ALIGN.CENTER
        
        # Value
        p2 = text_frame.add_paragraph()
        p2.text = str(value)
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = color
        p2.alignment = PP_ALIGN.CENTER
        
        # Description
        p3 = text_frame.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(107, 114, 128)
        p3.alignment = PP_ALIGN.CENTER

def create_risk_distribution_slide(prs, df, primary_blue, text_dark):
    """Create risk distribution slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Risk Distribution Overview"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    
    # Calculate statistics
    total_records = len(df)
    high_count = int((df["Priority"].str.upper() == "HIGH").sum()) if "Priority" in df.columns else 0
    med_count = int((df["Priority"].str.upper() == "MEDIUM").sum()) if "Priority" in df.columns else 0
    low_count = int((df["Priority"].str.upper() == "LOW").sum()) if "Priority" in df.columns else 0
    
    avg_risk = df["Risk_Score"].mean() if "Risk_Score" in df.columns else 0
    max_risk = df["Risk_Score"].max() if "Risk_Score" in df.columns else 0
    
    unique_locations = df["Location_Norm"].nunique() if "Location_Norm" in df.columns else 0
    unique_hazards = df["Type of Hazard"].nunique() if "Type of Hazard" in df.columns else 0
    
    stats_text = f"""Risk Assessment Statistics:

• Total Findings: {total_records}
• High Priority: {high_count} ({(high_count/total_records*100):.1f}%)
• Medium Priority: {med_count} ({(med_count/total_records*100):.1f}%)
• Low Priority: {low_count} ({(low_count/total_records*100):.1f}%)

• Average Risk Score: {avg_risk:.2f}
• Maximum Risk Score: {max_risk:.1f}
• Locations Assessed: {unique_locations}
• Hazard Types Identified: {unique_hazards}"""
    
    content.text = stats_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_dark

def create_priority_slide(prs, title_text, priority_content, color, text_dark, priority_df):
    """Create priority analysis slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = color
    
    content = slide.placeholders[1]
    
    # Clean and format content
    clean_content = clean_html_text(priority_content)
    
    # Add top findings if available
    if len(priority_df) > 0:
        top_findings = priority_df.nlargest(3, "Risk_Score") if "Risk_Score" in priority_df.columns else priority_df.head(3)
        findings_text = "\n\nTop Findings:\n"
        for i, (_, row) in enumerate(top_findings.iterrows(), 1):
            obs = str(row.get("Observation", ""))[:100] + "..." if len(str(row.get("Observation", ""))) > 100 else str(row.get("Observation", ""))
            findings_text += f"{i}. Risk Score: {row.get('Risk_Score', 0):.1f} - {obs}\n"
        clean_content += findings_text
    
    content.text = clean_content
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = text_dark

def create_location_analysis_slide(prs, df, primary_blue, text_dark):
    """Create location analysis slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Location Risk Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    
    if "Location_Norm" in df.columns and "Risk_Score" in df.columns:
        location_risks = df.groupby("Location_Norm")["Risk_Score"].agg(['sum', 'mean', 'count']).round(2)
        location_risks = location_risks.sort_values('sum', ascending=False)
        
        analysis_text = "Risk Assessment by Location:\n\n"
        for location, (total_risk, avg_risk, count) in location_risks.head(10).iterrows():
            analysis_text += f"• {location}: Total Risk: {total_risk}, Avg: {avg_risk}, Findings: {count}\n"
            
        content.text = analysis_text
    else:
        content.text = "Location analysis data not available"
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = text_dark

def create_hazard_analysis_slide(prs, df, primary_blue, text_dark):
    """Create hazard type analysis slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Hazard Type Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    
    if "Type of Hazard" in df.columns and "Risk_Score" in df.columns:
        hazard_risks = df.groupby("Type of Hazard")["Risk_Score"].agg(['sum', 'mean', 'count']).round(2)
        hazard_risks = hazard_risks.sort_values('sum', ascending=False)
        
        analysis_text = "Risk Assessment by Hazard Type:\n\n"
        for hazard, (total_risk, avg_risk, count) in hazard_risks.head(8).iterrows():
            analysis_text += f"• {hazard}: Total Risk: {total_risk}, Avg: {avg_risk}, Count: {count}\n"
            
        content.text = analysis_text
    else:
        content.text = "Hazard type analysis data not available"
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = text_dark

def create_top_findings_slide(prs, df, primary_blue, danger_red, text_dark):
    """Create top findings slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Top Risk Findings"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    
    if "Risk_Score" in df.columns:
        top_risks = df.nlargest(5, "Risk_Score")
        findings_text = ""
        
        for i, (_, row) in enumerate(top_risks.iterrows(), 1):
            risk_score = row.get("Risk_Score", 0)
            location = row.get("Location_Norm", "Unknown")
            hazard = row.get("Type of Hazard", "Unknown")
            priority = row.get("Priority", "Unknown")
            observation = str(row.get("Observation", ""))[:150] + "..." if len(str(row.get("Observation", ""))) > 150 else str(row.get("Observation", ""))
            
            findings_text += f"{i}. Risk Score: {risk_score}\n"
            findings_text += f"   Location: {location} | Hazard: {hazard} | Priority: {priority}\n"
            findings_text += f"   {observation}\n\n"
        
        content.text = findings_text
    else:
        content.text = "Risk findings data not available"
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = text_dark

def create_recommendations_slide(prs, df, kpis, primary_blue, text_dark):
    """Create recommendations and next steps slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Recommendations & Next Steps"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    content = slide.placeholders[1]
    
    recommendations = f"""Immediate Action Items:

• Address {kpis['high']} critical high-priority risks requiring immediate attention
• Develop action plans for {kpis['med']} medium-priority findings
• Schedule routine maintenance for {kpis['low']} low-priority items

Strategic Recommendations:

• Implement comprehensive risk monitoring system
• Establish regular safety audit schedules
• Enhance staff training on identified hazard types
• Review and update safety protocols based on findings
• Allocate resources for high-risk location improvements

Next Steps:

• Prioritize critical findings for immediate remediation
• Assign responsible teams for each risk category
• Set target completion dates for all action items
• Schedule follow-up assessments to measure progress
• Document lessons learned for future audits"""
    
    content.text = recommendations
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = text_dark
        paragraph.space_after = Pt(8)

def clean_html_text(html_text):
    """Remove HTML tags and clean text for PowerPoint"""
    import re
    if not html_text:
        return "Content not available"
    
    # Remove HTML tags
    clean = re.sub('<[^<]+?>', '', str(html_text))
    # Clean up common HTML entities
    clean = clean.replace('&nbsp;', ' ')
    clean = clean.replace('&amp;', '&')
    clean = clean.replace('&lt;', '<')
    clean = clean.replace('&gt;', '>')
    # Remove extra whitespace
    clean = re.sub(r'\s+', ' ', clean).strip()
    
    return clean
# ---------- HTML BUILDER ----------
def build_html(path, company, audit_date, kpis, charts, chart_summaries, exec_summary, priority_details, records, logo_base64):
    env = Environment(autoescape=select_autoescape(["html","xml"]))
    tpl = env.from_string(r"""
<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{{ company }} - Audit Risk Report</title>
  <script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
  <script src="https://cdnjs.cloudflare.com/ajax/libs/html2canvas/1.4.1/html2canvas.min.js"></script>
  <script src="https://cdnjs.cloudflare.com/ajax/libs/jspdf/2.5.1/jspdf.umd.min.js"></script>
  <link rel="stylesheet" href="https://cdn.datatables.net/1.13.6/css/jquery.dataTables.min.css">
  <script src="https://code.jquery.com/jquery-3.7.1.js"></script>
  <script src="https://cdn.datatables.net/1.13.6/js/jquery.dataTables.min.js"></script>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
  <style>
    :root {
      --primary-blue: #1e40af;
      --primary-light: #3b82f6;
      --danger-red: #dc2626;
      --warning-orange: #ea580c;
      --success-green: #059669;
      --dark-blue: #1e3a8a;
      --text-dark: #1f2937;
      --text-medium: #374151;
      --text-light: #6b7280;
      --border-color: #d1d5db;
      --bg-light: #f9fafb;
      --bg-white: #ffffff;
      --shadow-light: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
      --shadow-heavy: 0 10px 15px -3px rgba(0, 0, 0, 0.1);
      --border-radius: 12px;
      --transition: all 0.3s ease;
    }

    * { box-sizing: border-box; margin: 0; padding: 0; }

    body {
      font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
      background: linear-gradient(135deg, #1e40af 0%, #3b82f6 50%, #60a5fa 100%);
      min-height: 100vh;
      color: var(--text-dark);
      line-height: 1.6;
      padding: 20px;
    }

    .container {
      max-width: 1400px;
      margin: 0 auto;
      backdrop-filter: blur(20px);
      background: rgba(255, 255, 255, 0.95);
      border: 1px solid rgba(255, 255, 255, 0.3);
      border-radius: var(--border-radius);
      padding: 40px;
      box-shadow: var(--shadow-heavy);
      animation: fadeInUp 0.8s ease-out;
    }

    @keyframes fadeInUp {
      from { opacity: 0; transform: translateY(30px); }
      to { opacity: 1; transform: translateY(0); }
    }

    .header {
      text-align: center;
      margin-bottom: 40px;
      position: relative;
    }

    .logo-container {
      margin-bottom: 20px;
      display: flex;
      justify-content: center;
      align-items: center;
    }

    .company-logo {
      max-height: 120px;
      max-width: 300px;
      height: auto;
      filter: drop-shadow(0 4px 8px rgba(30, 64, 175, 0.2));
      border-radius: 8px;
      background: white;
      padding: 10px;
      box-shadow: var(--shadow-light);
    }

    .company-name {
      font-size: 3rem;
      font-weight: 700;
      color: var(--primary-blue);
      margin-bottom: 10px;
      text-shadow: 0 2px 4px rgba(30, 64, 175, 0.3);
      letter-spacing: -0.02em;
    }

    .audit-title {
      font-size: 1.5rem;
      color: var(--text-dark);
      font-weight: 500;
      margin-bottom: 8px;
    }

    .audit-date {
      font-size: 1.1rem;
      color: var(--text-medium);
      font-family: 'JetBrains Mono', monospace;
      font-weight: 500;
    }

    .kpi-grid {
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
      gap: 24px;
      margin-bottom: 50px;
    }

    .kpi-card {
      background: var(--bg-white);
      border-radius: var(--border-radius);
      padding: 32px 24px;
      text-align: center;
      cursor: pointer;
      transition: var(--transition);
      box-shadow: var(--shadow-light);
      border: 2px solid transparent;
      position: relative;
      overflow: hidden;
    }

    .kpi-card::before {
      content: '';
      position: absolute;
      top: 0;
      left: 0;
      right: 0;
      height: 4px;
      background: var(--primary-blue);
      transition: var(--transition);
    }

    .kpi-card.high::before { background: var(--danger-red); }
    .kpi-card.medium::before { background: var(--warning-orange); }
    .kpi-card.low::before { background: var(--success-green); }

    .kpi-card:hover {
      transform: translateY(-8px) scale(1.02);
      box-shadow: var(--shadow-heavy);
      border-color: var(--primary-light);
    }

    .kpi-label {
      font-size: 0.95rem;
      font-weight: 600;
      color: var(--text-medium);
      margin-bottom: 12px;
      text-transform: uppercase;
      letter-spacing: 0.5px;
    }

    .kpi-value {
      font-size: 3rem;
      font-weight: 700;
      color: var(--text-dark);
      margin-bottom: 8px;
    }

    .kpi-subtitle {
      font-size: 0.85rem;
      color: var(--text-light);
      font-weight: 500;
    }

    .download-section {
      text-align: center;
      margin-bottom: 50px;
      display: flex;
      gap: 20px;
      justify-content: center;
      flex-wrap: wrap;
    }

    .download-btn {
      background: var(--primary-blue);
      color: white;
      border: none;
      padding: 16px 32px;
      border-radius: 50px;
      font-size: 1.1rem;
      font-weight: 600;
      cursor: pointer;
      transition: var(--transition);
      box-shadow: var(--shadow-light);
      display: inline-flex;
      align-items: center;
      gap: 12px;
      min-width: 200px;
    }

    .download-btn:hover {
      transform: translateY(-2px);
      box-shadow: var(--shadow-heavy);
      background: var(--primary-light);
    }

    .download-btn:disabled {
      opacity: 0.6;
      cursor: not-allowed;
      transform: none;
    }

    .section {
      background: var(--bg-white);
      border-radius: var(--border-radius);
      padding: 32px;
      margin-bottom: 32px;
      box-shadow: var(--shadow-light);
      border: 2px solid var(--border-color);
      animation: slideInLeft 0.6s ease-out;
    }

    @keyframes slideInLeft {
      from { opacity: 0; transform: translateX(-30px); }
      to { opacity: 1; transform: translateX(0); }
    }

    .section-title {
      font-size: 2rem;
      font-weight: 600;
      color: var(--primary-blue);
      margin-bottom: 24px;
      position: relative;
      padding-bottom: 12px;
    }

    .section-title::after {
      content: '';
      position: absolute;
      bottom: 0;
      left: 0;
      width: 60px;
      height: 4px;
      background: var(--primary-blue);
      border-radius: 2px;
    }

    .exec-summary {
      font-size: 1.1rem;
      line-height: 1.8;
      color: var(--text-dark);
      background: var(--bg-light);
      padding: 24px;
      border-radius: 16px;
      border-left: 4px solid var(--primary-blue);
      margin-bottom: 16px;
      border: 1px solid var(--border-color);
    }

    .chart-container {
      background: var(--bg-white);
      border-radius: var(--border-radius);
      padding: 32px;
      margin-bottom: 40px;
      box-shadow: var(--shadow-light);
      border: 2px solid var(--border-color);
      transition: var(--transition);
    }

    .chart-container:hover {
      transform: translateY(-4px);
      box-shadow: var(--shadow-heavy);
      border-color: var(--primary-light);
    }

    .chart-title {
      font-size: 1.5rem;
      font-weight: 600;
      color: var(--primary-blue);
      margin-bottom: 24px;
      text-align: center;
    }

    .insights-title {
      font-size: 1.2rem;
      font-weight: 600;
      color: var(--primary-blue);
      margin: 32px 0 16px 0;
      display: flex;
      align-items: center;
      gap: 8px;
    }

    .insights-title::before {
      content: '💡';
      font-size: 1.1rem;
    }

    .insights-list {
      list-style: none;
      padding: 0;
    }

    .insights-list li {
      background: var(--bg-light);
      margin-bottom: 12px;
      padding: 16px 20px;
      border-radius: 12px;
      border-left: 4px solid var(--primary-blue);
      font-size: 0.95rem;
      color: var(--text-dark);
      transition: var(--transition);
      position: relative;
      border: 1px solid var(--border-color);
    }

    .insights-list li:hover {
      transform: translateX(8px);
      background: #e0f2fe;
      border-left-color: var(--primary-light);
    }

    .priority-navigation {
      display: grid;
      grid-template-columns: repeat(auto-fit, minmax(150px, 1fr));
      gap: 16px;
      margin-bottom: 24px;
    }

    .priority-btn {
      background: var(--bg-white);
      border: 2px solid var(--border-color);
      padding: 12px 20px;
      border-radius: 12px;
      font-weight: 500;
      cursor: pointer;
      transition: var(--transition);
      text-align: center;
      font-size: 0.95rem;
      color: var(--text-dark);
    }

    .priority-btn:hover {
      background: var(--primary-blue);
      color: white;
      transform: translateY(-2px);
      box-shadow: var(--shadow-light);
      border-color: var(--primary-blue);
    }

    .priority-btn.active {
      background: var(--primary-blue);
      color: white;
      box-shadow: var(--shadow-light);
      border-color: var(--primary-blue);
    }

    .priority-content {
      background: var(--bg-light);
      padding: 24px;
      border-radius: 16px;
      border: 1px solid var(--border-color);
      font-size: 1rem;
      line-height: 1.7;
      color: var(--text-dark);
    }

    .data-table-container {
      background: var(--bg-white);
      border-radius: var(--border-radius);
      padding: 24px;
      box-shadow: var(--shadow-light);
      overflow: hidden;
      border: 2px solid var(--border-color);
    }

    table.dataTable {
      background: transparent !important;
      border-radius: 12px !important;
      overflow: hidden !important;
      border-collapse: separate !important;
      border-spacing: 0 !important;
      width: 100% !important;
    }

    table.dataTable thead th {
      background: var(--primary-blue) !important;
      color: white !important;
      font-weight: 600 !important;
      padding: 16px 12px !important;
      border: none !important;
      font-size: 0.9rem !important;
      text-transform: uppercase !important;
      letter-spacing: 0.5px !important;
    }

    table.dataTable tbody td {
      padding: 14px 12px !important;
      border-bottom: 1px solid var(--border-color) !important;
      font-size: 0.9rem !important;
      color: var(--text-dark) !important;
      vertical-align: top !important;
    }

    table.dataTable tbody tr:hover {
      background: var(--bg-light) !important;
    }

    .dataTables_wrapper .dataTables_length,
    .dataTables_wrapper .dataTables_filter,
    .dataTables_wrapper .dataTables_info,
    .dataTables_wrapper .dataTables_paginate {
      color: var(--text-medium) !important;
      font-weight: 500 !important;
    }

    .dataTables_wrapper .dataTables_filter input {
      background: var(--bg-white);
      border: 2px solid var(--border-color);
      border-radius: 8px;
      padding: 8px 12px;
      margin-left: 8px;
      transition: var(--transition);
      color: var(--text-dark);
    }

    .dataTables_wrapper .dataTables_filter input:focus {
      outline: none;
      border-color: var(--primary-blue);
      box-shadow: 0 0 0 3px rgba(30, 64, 175, 0.1);
    }

    .status-badge {
      padding: 6px 12px;
      border-radius: 20px;
      font-size: 0.8rem;
      font-weight: 600;
      text-transform: uppercase;
      letter-spacing: 0.5px;
    }

    .status-high {
      background: var(--danger-red);
      color: white;
    }

    .status-medium {
      background: var(--warning-orange);
      color: white;
    }

    .status-low {
      background: var(--success-green);
      color: white;
    }

    /* Print-specific styles for PDF generation */
    @media print {
      body {
        background: white !important;
        padding: 0 !important;
      }
      
      .container {
        background: white !important;
        box-shadow: none !important;
        border: none !important;
        backdrop-filter: none !important;
        max-width: none !important;
        padding: 20px !important;
        margin: 0 !important;
      }
      
      .floating-elements,
      .scroll-indicator {
        display: none !important;
      }
      
      .download-section {
        display: none !important;
      }
      
      .chart-container,
      .section {
        break-inside: avoid;
        page-break-inside: avoid;
      }
      
      .kpi-grid {
        break-inside: avoid;
      }
      
      .chart-container {
        margin-bottom: 30px !important;
      }
    }

    /* Floating background elements */
    .floating-elements {
      position: fixed;
      top: 0;
      left: 0;
      width: 100%;
      height: 100%;
      pointer-events: none;
      z-index: -1;
    }

    .floating-circle {
      position: absolute;
      border-radius: 50%;
      background: rgba(255, 255, 255, 0.1);
      animation: float 6s ease-in-out infinite;
    }

    .floating-circle:nth-child(1) {
      width: 80px;
      height: 80px;
      top: 10%;
      left: 10%;
      animation-delay: 0s;
    }

    .floating-circle:nth-child(2) {
      width: 120px;
      height: 120px;
      top: 60%;
      right: 10%;
      animation-delay: 2s;
    }

    .floating-circle:nth-child(3) {
      width: 60px;
      height: 60px;
      bottom: 20%;
      left: 20%;
      animation-delay: 4s;
    }

    @keyframes float {
      0%, 100% { transform: translateY(0px) rotate(0deg); }
      50% { transform: translateY(-20px) rotate(180deg); }
    }

    .scroll-indicator {
      position: fixed;
      top: 0;
      left: 0;
      width: 100%;
      height: 4px;
      background: rgba(255, 255, 255, 0.3);
      z-index: 1000;
    }

    .scroll-progress {
      height: 100%;
      background: var(--primary-blue);
      width: 0%;
      transition: width 0.3s ease;
    }

    /* Enhanced KPI styling */
    .kpi-card.high {
      border-left: 6px solid var(--danger-red);
    }

    .kpi-card.medium {
      border-left: 6px solid var(--warning-orange);
    }

    .kpi-card.low {
      border-left: 6px solid var(--success-green);
    }

    .kpi-card:nth-child(4) {
      border-left: 6px solid var(--primary-blue);
    }

    .chart-container .plotly-graph-div {
      border-radius: 8px;
    }

    .insights-list li,
    .exec-summary,
    .priority-content {
      font-weight: 400;
    }

    .exec-summary p {
      margin-bottom: 16px;
      font-weight: 400;
    }

    .progress-container {
      background: var(--bg-light);
      border-radius: 8px;
      padding: 20px;
      margin-top: 20px;
      border: 1px solid var(--border-color);
      text-align: center;
    }

    .progress-bar {
      width: 100%;
      height: 8px;
      background: var(--border-color);
      border-radius: 4px;
      overflow: hidden;
      margin: 10px 0;
    }

    .progress-fill {
      height: 100%;
      background: var(--primary-blue);
      width: 0%;
      transition: width 0.3s ease;
      border-radius: 4px;
    }

    @media (max-width: 768px) {
      body { padding: 10px; }
      .container { padding: 20px; }
      .company-name { font-size: 2rem; }
      .company-logo { max-height: 80px; }
      .kpi-grid { grid-template-columns: 1fr; }
      .chart-container { padding: 20px; }
      .download-section { flex-direction: column; align-items: center; }
    }

    /* Enhanced animations */
    .kpi-card {
      animation: bounceIn 0.6s ease-out;
      animation-fill-mode: both;
    }

    .kpi-card:nth-child(1) { animation-delay: 0.1s; }
    .kpi-card:nth-child(2) { animation-delay: 0.2s; }
    .kpi-card:nth-child(3) { animation-delay: 0.3s; }
    .kpi-card:nth-child(4) { animation-delay: 0.4s; }

    @keyframes bounceIn {
      0% { opacity: 0; transform: scale(0.3) translateY(50px); }
      50% { opacity: 1; transform: scale(1.05) translateY(-10px); }
      70% { transform: scale(0.9) translateY(0px); }
      100% { opacity: 1; transform: scale(1) translateY(0px); }
    }

    .chart-container {
      animation: slideInRight 0.8s ease-out;
      animation-fill-mode: both;
    }

    .chart-container:nth-child(even) {
      animation: slideInLeft 0.8s ease-out;
      animation-fill-mode: both;
    }

    @keyframes slideInRight {
      from { opacity: 0; transform: translateX(50px); }
      to { opacity: 1; transform: translateX(0); }
    }

    @keyframes slideInLeft {
      from { opacity: 0; transform: translateX(-50px); }
      to { opacity: 1; transform: translateX(0); }
    }
  </style>
</head>
<body>
  <div class="floating-elements">
    <div class="floating-circle"></div>
    <div class="floating-circle"></div>
    <div class="floating-circle"></div>
  </div>
  
  <div class="scroll-indicator">
    <div class="scroll-progress" id="scrollProgress"></div>
  </div>

  <div class="container" id="reportContainer">
    <div class="header">
      {% if logo_base64 %}
      <div class="logo-container">
        <img src="{{ logo_base64 }}" alt="{{ company }} Logo" class="company-logo">
      </div>
      {% endif %}
      <h1 class="company-name">{{ company }}</h1>
      <h2 class="audit-title">Comprehensive Risk Audit Report</h2>
      <p class="audit-date">Generated on {{ audit_date }}</p>
    </div>

    <div class="kpi-grid">
      <div class="kpi-card high" onclick="showPriority('HIGH')">
        <div class="kpi-label">Critical Risks</div>
        <div class="kpi-value">{{ kpis.high }}</div>
        <div class="kpi-subtitle">Immediate Action Required</div>
      </div>
      <div class="kpi-card medium" onclick="showPriority('MEDIUM')">
        <div class="kpi-label">Medium Risks</div>
        <div class="kpi-value">{{ kpis.med }}</div>
        <div class="kpi-subtitle">Monitor & Plan</div>
      </div>
      <div class="kpi-card low" onclick="showPriority('LOW')">
        <div class="kpi-label">Low Risks</div>
        <div class="kpi-value">{{ kpis.low }}</div>
        <div class="kpi-subtitle">Routine Maintenance</div>
      </div>
      <div class="kpi-card" onclick="showPriority('ALL')">
        <div class="kpi-label">Total Findings</div>
        <div class="kpi-value">{{ kpis.total }}</div>
        <div class="kpi-subtitle">Complete Assessment</div>
      </div>
    </div>

    <div class="download-section">
      <button class="download-btn" onclick="downloadComprehensivePDF()">
        <span>📄</span>
        Download Complete PDF Report
      </button>
      <button class="download-btn" onclick="printReport()" style="background: var(--success-green);">
        <span>🖨️</span>
        Print Report
      </button>
    </div>

    <div class="section">
      <h2 class="section-title">Executive Summary</h2>
      <div class="exec-summary">{{ exec_summary|safe }}</div>
    </div>

    {% for name, chart_html in charts.items() %}
      <div class="chart-container" id="chart-{{ loop.index }}">
        <h3 class="chart-title">{{ name }}</h3>
        <div class="chart-content">{{ chart_html|safe }}</div>
        <h4 class="insights-title">Key Insights</h4>
        <ul class="insights-list">
          {% for line in chart_summaries[name].splitlines() %}
            {% if line.strip() %}<li>{{ line.strip().lstrip('•').lstrip('-').strip() }}</li>{% endif %}
          {% endfor %}
        </ul>
      </div>
    {% endfor %}

    <div class="section">
      <h2 class="section-title">Priority Analysis</h2>
      <div class="priority-navigation">
        <button class="priority-btn active" onclick="showPriority('ALL')" id="btn-ALL">All Findings</button>
        <button class="priority-btn" onclick="showPriority('HIGH')" id="btn-HIGH">High Priority</button>
        <button class="priority-btn" onclick="showPriority('MEDIUM')" id="btn-MEDIUM">Medium Priority</button>
        <button class="priority-btn" onclick="showPriority('LOW')" id="btn-LOW">Low Priority</button>
      </div>
      <div id="priorityContent" class="priority-content">{{ priority_details['ALL']|safe }}</div>
    </div>

    <div class="section">
      <h2 class="section-title">Complete Audit Dataset</h2>
      <div class="data-table-container">
        <table id="auditTable" class="display" style="width:100%">
          <thead>
            <tr>
              {% for col in records[0].keys() %}
                <th>{{ col }}</th>
              {% endfor %}
            </tr>
          </thead>
          <tbody>
            {% for row in records %}
            <tr>
              {% for col in records[0].keys() %}
                <td>
                  {% if col == 'Priority' %}
                    <span class="status-badge status-{{ row[col]|lower }}">{{ row[col] }}</span>
                  {% else %}
                    {{ row[col] }}
                  {% endif %}
                </td>
              {% endfor %}
            </tr>
            {% endfor %}
          </tbody>
        </table>
      </div>
    </div>
  </div>

  <script>
    // Store chart data and summaries for PDF generation
    const chartData = {{ charts|tojson }};
    const chartSummaries = {{ chart_summaries|tojson }};
    const kpiData = {{ kpis|tojson }};
    const priorityData = {{ priority_details|tojson }};
    const auditRecords = {{ records|tojson }};
    const companyName = "{{ company }}";
    const auditDate = "{{ audit_date }}";
    const execSummary = `{{ exec_summary|striptags|safe }}`;

    // Initialize DataTable
    $(document).ready(function() {
      $('#auditTable').DataTable({
        pageLength: 25,
        responsive: true,
        dom: 'frtip',
        order: [[0, 'asc']],
        columnDefs: [
          { targets: '_all', className: 'dt-center' }
        ]
      });
    });

    // Scroll progress indicator
    window.addEventListener('scroll', function() {
      const scrollTop = window.pageYOffset;
      const docHeight = document.body.offsetHeight - window.innerHeight;
      const scrollPercent = (scrollTop / docHeight) * 100;
      document.getElementById('scrollProgress').style.width = scrollPercent + '%';
    });

    // Enhanced PDF generation with all content and charts
    async function downloadComprehensivePDF() {
      const downloadBtn = document.querySelector('.download-btn');
      const originalText = downloadBtn.innerHTML;
      
      // Show loading state
      downloadBtn.innerHTML = '<span>⏳</span> Generating Comprehensive PDF...';
      downloadBtn.disabled = true;

      try {
        // Create progress indicator
        const progressDiv = document.createElement('div');
        progressDiv.className = 'progress-container';
        progressDiv.innerHTML = `
          <h4>Generating PDF Report...</h4>
          <div class="progress-bar">
            <div class="progress-fill" id="progressFill"></div>
          </div>
          <p id="progressText">Initializing...</p>
        `;
        downloadBtn.parentNode.appendChild(progressDiv);

        const updateProgress = (percent, text) => {
          document.getElementById('progressFill').style.width = percent + '%';
          document.getElementById('progressText').textContent = text;
        };

        updateProgress(10, 'Creating PDF document...');

        const { jsPDF } = window.jspdf;
        const doc = new jsPDF('p', 'mm', 'a4');
        const pageWidth = doc.internal.pageSize.getWidth();
        const pageHeight = doc.internal.pageSize.getHeight();
        let yPosition = 20;

        // Helper function to add new page if needed
        const checkNewPage = (requiredSpace = 20) => {
          if (yPosition + requiredSpace > pageHeight - 20) {
            doc.addPage();
            yPosition = 20;
            return true;
          }
          return false;
        };

        // Helper function to add wrapped text
        const addWrappedText = (text, x, maxWidth, fontSize = 10) => {
          doc.setFontSize(fontSize);
          const lines = doc.splitTextToSize(text, maxWidth);
          lines.forEach(line => {
            checkNewPage();
            doc.text(line, x, yPosition);
            yPosition += fontSize * 0.4;
          });
          return lines.length;
        };

        updateProgress(20, 'Adding header and logo...');

        // Add logo if available
        {% if logo_base64 %}
        try {
          doc.addImage("{{ logo_base64 }}", 'PNG', pageWidth/2 - 25, yPosition, 50, 20);
          yPosition += 30;
        } catch (e) {
          console.log('Logo embedding failed, continuing without logo');
        }
        {% endif %}

        // Add header
        doc.setFontSize(24);
        doc.setTextColor(30, 64, 175);
        doc.text(companyName, pageWidth/2, yPosition, { align: 'center' });
        yPosition += 12;
        
        doc.setFontSize(16);
        doc.setTextColor(55, 65, 81);
        doc.text('Comprehensive Risk Audit Report', pageWidth/2, yPosition, { align: 'center' });
        yPosition += 10;
        
        doc.setFontSize(12);
        doc.text(`Generated on ${auditDate}`, pageWidth/2, yPosition, { align: 'center' });
        yPosition += 20;

        updateProgress(30, 'Adding KPI summary...');

        // Add KPIs section
        checkNewPage(40);
        doc.setFontSize(18);
        doc.setTextColor(30, 64, 175);
        doc.text('Key Performance Indicators', 20, yPosition);
        yPosition += 15;

        doc.setFontSize(12);
        doc.setTextColor(0, 0, 0);
        
        // KPI table
        const kpiTable = [
          ['Critical Risks (High Priority)', kpiData.high, 'Immediate Action Required'],
          ['Medium Risks', kpiData.med, 'Monitor & Plan'],
          ['Low Risks', kpiData.low, 'Routine Maintenance'],
          ['Total Findings', kpiData.total, 'Complete Assessment']
        ];

        kpiTable.forEach(([label, value, desc]) => {
          checkNewPage();
          doc.setFontSize(11);
          doc.setTextColor(30, 64, 175);
          doc.text(`${label}:`, 20, yPosition);
          doc.setFontSize(14);
          doc.setTextColor(0, 0, 0);
          doc.text(String(value), 80, yPosition);
          doc.setFontSize(9);
          doc.setTextColor(100, 100, 100);
          doc.text(`(${desc})`, 90, yPosition);
          yPosition += 8;
        });
        yPosition += 10;

        updateProgress(40, 'Adding executive summary...');

        // Add Executive Summary
        checkNewPage(30);
        doc.setFontSize(18);
        doc.setTextColor(30, 64, 175);
        doc.text('Executive Summary', 20, yPosition);
        yPosition += 12;

        doc.setTextColor(0, 0, 0);
        addWrappedText(execSummary, 20, pageWidth - 40, 10);
        yPosition += 15;

        updateProgress(50, 'Adding priority analysis...');

        // Add Priority Analysis
        for (const [priority, content] of Object.entries(priorityData)) {
          checkNewPage(25);
          doc.setFontSize(16);
          doc.setTextColor(30, 64, 175);
          doc.text(`${priority} Priority Analysis`, 20, yPosition);
          yPosition += 10;

          doc.setTextColor(0, 0, 0);
          const cleanContent = content.replace(/<[^>]*>/g, '').replace(/&[^;]+;/g, ' ');
          addWrappedText(cleanContent, 20, pageWidth - 40, 9);
          yPosition += 10;
        }

        updateProgress(60, 'Adding charts and insights...');

        // Add Chart Summaries
        const chartNames = Object.keys(chartSummaries);
        for (let i = 0; i < chartNames.length; i++) {
          const chartName = chartNames[i];
          const insights = chartSummaries[chartName];

          checkNewPage(30);
          doc.setFontSize(16);
          doc.setTextColor(30, 64, 175);
          doc.text(`Chart Analysis: ${chartName}`, 20, yPosition);
          yPosition += 12;

          doc.setFontSize(10);
          doc.setTextColor(0, 0, 0);
          
          // Add insights
          const insightLines = insights.split('\n').filter(line => line.trim());
          insightLines.forEach(line => {
            const cleanLine = line.replace(/^[•\-\*]\s*/, '• ');
            checkNewPage();
            doc.text(cleanLine, 20, yPosition);
            yPosition += 6;
          });
          yPosition += 8;

          updateProgress(60 + (i / chartNames.length) * 20, `Processing chart ${i + 1} of ${chartNames.length}...`);
        }

        updateProgress(80, 'Adding detailed findings...');

        // Add Top Risk Findings Table
        checkNewPage(40);
        doc.setFontSize(18);
        doc.setTextColor(30, 64, 175);
        doc.text('Detailed Risk Findings', 20, yPosition);
        yPosition += 15;

        // Table headers
        doc.setFontSize(8);
        doc.setTextColor(255, 255, 255);
        doc.setFillColor(30, 64, 175);
        doc.rect(20, yPosition - 4, pageWidth - 40, 8, 'F');
        
        doc.text('Sr.', 22, yPosition);
        doc.text('Priority', 35, yPosition);
        doc.text('Risk Score', 55, yPosition);
        doc.text('Location', 75, yPosition);
        doc.text('Hazard Type', 105, yPosition);
        doc.text('Observation', 135, yPosition);
        yPosition += 8;

        // Add top 30 records with details
        doc.setTextColor(0, 0, 0);
        const topRecords = auditRecords.slice(0, 30);
        
        topRecords.forEach((record, index) => {
          checkNewPage(12);
          
          doc.setFontSize(7);
          doc.text(String(record['Sr. No'] || index + 1), 22, yPosition);
          
          // Color-code priority
          const priority = String(record['Priority'] || 'N/A').toUpperCase();
          if (priority === 'HIGH') doc.setTextColor(220, 38, 38);
          else if (priority === 'MEDIUM') doc.setTextColor(234, 88, 12);
          else if (priority === 'LOW') doc.setTextColor(5, 150, 105);
          else doc.setTextColor(0, 0, 0);
          
          doc.text(priority, 35, yPosition);
          doc.setTextColor(0, 0, 0);
          
          doc.text(String(record['Risk_Score'] || '0'), 55, yPosition);
          doc.text(String(record['Location_Norm'] || 'N/A').substring(0, 12), 75, yPosition);
          doc.text(String(record['Type of Hazard'] || 'N/A').substring(0, 15), 105, yPosition);
          doc.text(String(record['Observation'] || 'N/A').substring(0, 25), 135, yPosition);
          
          yPosition += 6;
          
          // Add recommendation if available
          const recommendation = String(record['Recommendation'] || '').trim();
          if (recommendation && recommendation !== 'N/A' && recommendation !== '') {
            checkNewPage(6);
            doc.setFontSize(6);
            doc.setTextColor(100, 100, 100);
            doc.text(`Rec: ${recommendation.substring(0, 60)}...`, 35, yPosition);
            yPosition += 4;
            doc.setTextColor(0, 0, 0);
          }
          
          updateProgress(80 + (index / topRecords.length) * 15, `Adding finding ${index + 1} of ${topRecords.length}...`);
        });

        updateProgress(95, 'Adding summary statistics...');

        // Add summary statistics
        checkNewPage(30);
        doc.setFontSize(16);
        doc.setTextColor(30, 64, 175);
        doc.text('Risk Assessment Summary', 20, yPosition);
        yPosition += 12;

        doc.setFontSize(10);
        doc.setTextColor(0, 0, 0);

        // Calculate additional statistics
        const totalRiskScore = auditRecords.reduce((sum, record) => sum + (record['Risk_Score'] || 0), 0);
        const avgRiskScore = (totalRiskScore / auditRecords.length).toFixed(2);
        const highRiskPercent = ((kpiData.high / kpiData.total) * 100).toFixed(1);
        const mediumRiskPercent = ((kpiData.med / kpiData.total) * 100).toFixed(1);
        const lowRiskPercent = ((kpiData.low / kpiData.total) * 100).toFixed(1);

        const summaryStats = [
          `Total Risk Score: ${totalRiskScore.toFixed(1)}`,
          `Average Risk Score: ${avgRiskScore}`,
          `High Priority: ${kpiData.high} findings (${highRiskPercent}%)`,
          `Medium Priority: ${kpiData.med} findings (${mediumRiskPercent}%)`,
          `Low Priority: ${kpiData.low} findings (${lowRiskPercent}%)`,
          `Report Generated: ${auditDate}`,
          `Total Locations Assessed: ${new Set(auditRecords.map(r => r['Location_Norm'])).size}`,
          `Hazard Types Identified: ${new Set(auditRecords.map(r => r['Type of Hazard'])).size}`
        ];

        summaryStats.forEach(stat => {
          checkNewPage();
          doc.text(`• ${stat}`, 20, yPosition);
          yPosition += 6;
        });

        updateProgress(100, 'Finalizing PDF...');

        // Add footer to each page
        const totalPages = doc.internal.getNumberOfPages();
        for (let i = 1; i <= totalPages; i++) {
          doc.setPage(i);
          doc.setFontSize(8);
          doc.setTextColor(150, 150, 150);
          doc.text(`${companyName} - Audit Report`, 20, pageHeight - 10);
          doc.text(`Page ${i} of ${totalPages}`, pageWidth - 40, pageHeight - 10);
          doc.text(`Generated: ${auditDate}`, pageWidth/2, pageHeight - 10, { align: 'center' });
        }

        // Save PDF
        const filename = `Chola_MS_Risk_Services_Audit_Report_${auditDate.replace(/-/g, '_')}.pdf`;
        doc.save(filename);
        
        // Remove progress indicator
        progressDiv.remove();
        
        // Restore button
        downloadBtn.innerHTML = originalText;
        downloadBtn.disabled = false;
        
        console.log('Comprehensive PDF generated successfully');
        
      } catch (error) {
        console.error('PDF generation error:', error);
        
        // Remove progress indicator if it exists
        const progressDiv = document.querySelector('.progress-container');
        if (progressDiv) progressDiv.remove();
        
        downloadBtn.innerHTML = originalText;
        downloadBtn.disabled = false;
        
        alert('PDF generation failed. Error: ' + error.message + '\nPlease try the print option instead.');
      }
    }

    // Simple print function as backup
    function printReport() {
      // Hide download buttons for printing
      document.querySelectorAll('.download-btn').forEach(btn => {
        btn.style.display = 'none';
      });
      
      // Trigger print
      window.print();
      
      // Restore download buttons after print dialog
      setTimeout(() => {
        document.querySelectorAll('.download-btn').forEach(btn => {
          btn.style.display = 'inline-flex';
        });
      }, 1000);
    }

    // Priority filtering with smooth transitions
    function showPriority(priority) {
      // Update button states
      document.querySelectorAll('.priority-btn').forEach(btn => {
        btn.classList.remove('active');
      });
      document.getElementById('btn-' + priority).classList.add('active');

      // Update content with fade effect
      const contentDiv = document.getElementById("priorityContent");
      contentDiv.style.opacity = '0';
      contentDiv.style.transform = 'translateY(20px)';

      setTimeout(() => {
        contentDiv.innerHTML = priorityData[priority];
        contentDiv.style.opacity = '1';
        contentDiv.style.transform = 'translateY(0)';
      }, 200);

      // Filter DataTable
      const table = $('#auditTable').DataTable();
      if (priority === 'ALL') {
        table.search('').draw();
      } else {
        table.search(priority).draw();
      }
    }

    // Add smooth scroll behavior
    document.documentElement.style.scrollBehavior = 'smooth';

    // Enhanced KPI card interactions
    document.querySelectorAll('.kpi-card').forEach(card => {
      card.addEventListener('mouseenter', function() {
        if (window.innerWidth > 768) {
          this.style.transform = 'translateY(-8px) scale(1.02)';
        }
      });
      
      card.addEventListener('mouseleave', function() {
        this.style.transform = 'translateY(0) scale(1)';
      });
    });

    // Intersection Observer for animations
    const observerOptions = {
      threshold: 0.1,
      rootMargin: '0px 0px -50px 0px'
    };

    const observer = new IntersectionObserver((entries) => {
      entries.forEach(entry => {
        if (entry.isIntersecting) {
          entry.target.style.opacity = '1';
          entry.target.style.transform = 'translateY(0)';
        }
      });
    }, observerOptions);

    // Observe chart containers for scroll animations
    document.querySelectorAll('.chart-container').forEach((chart, index) => {
      chart.style.opacity = '0';
      chart.style.transform = 'translateY(30px)';
      chart.style.transition = 'opacity 0.6s ease, transform 0.6s ease';
      chart.style.animationDelay = `${index * 0.1}s`;
      observer.observe(chart);
    });

    // Enhanced table styling after DataTable initialization
    setTimeout(() => {
      // Style pagination buttons
      document.querySelectorAll('.paginate_button').forEach(btn => {
        btn.style.background = 'rgba(255, 255, 255, 0.9)';
        btn.style.border = '1px solid #d1d5db';
        btn.style.borderRadius = '8px';
        btn.style.margin = '0 2px';
        btn.style.transition = 'all 0.2s ease';
        btn.style.color = '#374151';
      });

      // Add hover effects to pagination
      document.querySelectorAll('.paginate_button').forEach(btn => {
        btn.addEventListener('mouseenter', function() {
          this.style.background = '#1e40af';
          this.style.color = 'white';
        });
        btn.addEventListener('mouseleave', function() {
          this.style.background = 'rgba(255, 255, 255, 0.9)';
          this.style.color = '#374151';
        });
      });
    }, 100);

    // Company name emphasis
    document.addEventListener('DOMContentLoaded', function() {
      const companyName = document.querySelector('.company-name');
      if (companyName) {
        companyName.style.fontWeight = '700';
        companyName.style.letterSpacing = '-0.025em';
      }
    });

    // Auto-resize charts for better display
    window.addEventListener('resize', function() {
      const plots = document.querySelectorAll('.plotly-graph-div');
      plots.forEach(plot => {
        if (window.Plotly) {
          window.Plotly.Plots.resize(plot);
        }
      });
    });

    // Enhanced error handling for PDF generation
    window.addEventListener('error', function(e) {
      console.error('Runtime error:', e.error);
      const downloadBtn = document.querySelector('.download-btn:disabled');
      if (downloadBtn) {
        downloadBtn.innerHTML = '<span>📄</span> Download Complete PDF Report';
        downloadBtn.disabled = false;
      }
    });
  </script>
</body>
</html>
""")
    rendered = tpl.render(
        company=company,
        audit_date=audit_date,
        kpis=kpis,
        charts=charts,
        chart_summaries=chart_summaries,
        exec_summary=exec_summary or "No summary available.",
        priority_details=priority_details,
        records=records,
        logo_base64=logo_base64
    )
    with open(path,"w",encoding="utf-8") as f: f.write(rendered)
    print("✅ Enhanced PDF-optimized interactive report with logo generated:", path)

# ---------- MAIN ----------
# First, add all the PPT generation functions BEFORE the main() function
# (Insert all the PPT code from my previous response here)

# Then, replace your main() function with this corrected version:

def main():
    df = load_data(DATA_FILE)
    audit_date = AUDIT_DATE or datetime.datetime.now(datetime.timezone.utc).strftime("%Y-%m-%d")
    df = prepare_records(df)
    kpis = compute_kpis(df)

    # Load and encode company logo
    logo_base64 = encode_image_to_base64(COMPANY_LOGO)
    if logo_base64:
        print("✅ Company logo loaded successfully")
    else:
        print("ℹ️ Continuing without logo")

    charts, summaries = {}, {}

    # Configure Plotly for better PDF rendering
    plotly_config = {
        'displayModeBar': False,
        'staticPlot': False,
        'responsive': True,
        'toImageButtonOptions': {
            'format': 'png',
            'filename': 'chart',
            'height': 500,
            'width': 800,
            'scale': 2
        }
    }

    print("📊 Generating charts...")

    # Hazard Risk Breakdown
    fig1 = px.bar(df, x="Type of Hazard", y="Risk_Score", color="Priority", text="Risk_Score", 
                  title="Hazard Risk Breakdown",
                  color_discrete_map={'HIGH': '#dc2626', 'MEDIUM': '#ea580c', 'LOW': '#059669'})
    fig1.update_layout(
        plot_bgcolor='white',
        paper_bgcolor='white',
        font={'color': '#1f2937', 'size': 12},
        title={'font': {'size': 16, 'color': '#1e40af'}},
        xaxis={'title': {'font': {'color': '#374151'}}},
        yaxis={'title': {'font': {'color': '#374151'}}},
        showlegend=True,
        margin=dict(l=50, r=50, t=60, b=50)
    )
    fig1.update_traces(textposition='outside')
    charts["Hazard Risk Breakdown"] = fig1.to_html(full_html=False, include_plotlyjs="cdn", config=plotly_config)
    summaries["Hazard Risk Breakdown"] = generate_chart_summary("Hazard Risk Breakdown", df)

    # Location vs Hazard Analysis
    fig2 = px.density_heatmap(df, x="Location_Norm", y="Type of Hazard", z="Risk_Score", 
                              title="Location vs Hazard Analysis", color_continuous_scale="Blues")
    fig2.update_layout(
        plot_bgcolor='white',
        paper_bgcolor='white',
        font={'color': '#1f2937', 'size': 12},
        title={'font': {'size': 16, 'color': '#1e40af'}},
        xaxis={'title': {'font': {'color': '#374151'}}},
        yaxis={'title': {'font': {'color': '#374151'}}},
        margin=dict(l=50, r=50, t=60, b=50)
    )
    charts["Location vs Hazard Analysis"] = fig2.to_html(full_html=False, include_plotlyjs=False, config=plotly_config)
    summaries["Location vs Hazard Analysis"] = generate_chart_summary("Location vs Hazard Analysis", df)

    # Risk Contribution by Hazard Type
    risk_by_hazard = df.groupby("Type of Hazard")["Risk_Score"].sum().reset_index()
    fig3 = px.bar(risk_by_hazard, x="Type of Hazard", y="Risk_Score", text="Risk_Score", 
                  title="Risk Contribution by Hazard Type")
    fig3.update_layout(
        plot_bgcolor='white',
        paper_bgcolor='white',
        font={'color': '#1f2937', 'size': 12},
        title={'font': {'size': 16, 'color': '#1e40af'}},
        xaxis={'title': {'font': {'color': '#374151'}}},
        yaxis={'title': {'font': {'color': '#374151'}}},
        margin=dict(l=50, r=50, t=60, b=50)
    )
    fig3.update_traces(marker_color='#1e40af', textposition='outside')
    charts["Risk Contribution by Hazard Type"] = fig3.to_html(full_html=False, include_plotlyjs=False, config=plotly_config)
    summaries["Risk Contribution by Hazard Type"] = generate_chart_summary("Risk Contribution by Hazard Type", df)

    # Hierarchical Risk Distribution
    fig4 = px.sunburst(df, path=["Location_Norm","Type of Hazard","Priority"], values="Risk_Score", 
                       title="Hierarchical Risk Distribution",
                       color_discrete_sequence=['#1e40af', '#3b82f6', '#60a5fa', '#93c5fd'])
    fig4.update_layout(
        plot_bgcolor='white',
        paper_bgcolor='white',
        font={'color': '#1f2937', 'size': 12},
        title={'font': {'size': 16, 'color': '#1e40af'}},
        margin=dict(l=50, r=50, t=60, b=50)
    )
    charts["Hierarchical Risk Distribution"] = fig4.to_html(full_html=False, include_plotlyjs=False, config=plotly_config)
    summaries["Hierarchical Risk Distribution"] = generate_chart_summary("Hierarchical Risk Distribution", df)

    # Risk vs Priority Correlation
    fig5 = px.scatter(df, x="Priority_Score", y="Risk_Score", size="Risk_Score", color="Type of Hazard", 
                      hover_name="Observation", title="Risk vs Priority Correlation")
    fig5.update_layout(
        plot_bgcolor='white',
        paper_bgcolor='white',
        font={'color': '#1f2937', 'size': 12},
        title={'font': {'size': 16, 'color': '#1e40af'}},
        xaxis={'title': {'font': {'color': '#374151'}}},
        yaxis={'title': {'font': {'color': '#374151'}}},
        margin=dict(l=50, r=50, t=60, b=50)
    )
    charts["Risk vs Priority Correlation"] = fig5.to_html(full_html=False, include_plotlyjs=False, config=plotly_config)
    summaries["Risk vs Priority Correlation"] = generate_chart_summary("Risk vs Priority Correlation", df)

    # Total Risk Assessment by Location
    risk_by_location = df.groupby("Location_Norm")["Risk_Score"].sum().reset_index()
    fig6 = px.bar(risk_by_location, x="Location_Norm", y="Risk_Score", text="Risk_Score", 
                  title="Total Risk Assessment by Location")
    fig6.update_layout(
        plot_bgcolor='white',
        paper_bgcolor='white',
        font={'color': '#1f2937', 'size': 12},
        title={'font': {'size': 16, 'color': '#1e40af'}},
        xaxis={'title': {'font': {'color': '#374151'}}},
        yaxis={'title': {'font': {'color': '#374151'}}},
        margin=dict(l=50, r=50, t=60, b=50)
    )
    fig6.update_traces(marker_color='#1e40af', textposition='outside')
    charts["Total Risk Assessment by Location"] = fig6.to_html(full_html=False, include_plotlyjs=False, config=plotly_config)
    summaries["Total Risk Assessment by Location"] = generate_chart_summary("Total Risk Assessment by Location", df)

    print("🤖 Generating AI summaries...")

    # Executive Summary
    top_examples = df.sort_values("Risk_Score", ascending=False).head(8).to_dict("records")
    exec_summary = generate_exec_summary(kpis, top_examples)

    # Priority details
    priority_details = {p: generate_priority_details(df, p) for p in ["HIGH","MEDIUM","LOW","ALL"]}

    print("🔨 Building HTML report...")
    build_html(OUTPUT_HTML, COMPANY_NAME, audit_date, kpis, charts, summaries, exec_summary, priority_details, df.to_dict("records"), logo_base64)

    # ADD PPT GENERATION HERE
    print("📊 Generating PowerPoint presentation...")
    ppt_path = generate_ppt_report(COMPANY_NAME, audit_date, kpis, exec_summary, priority_details, df, COMPANY_LOGO)
    if ppt_path:
        print(f"✅ PowerPoint presentation created: {ppt_path}")
    else:
        print("❌ PowerPoint generation failed")

# This function should be defined OUTSIDE and BEFORE main()
def generate_ppt_report(company_name, audit_date, kpis, exec_summary, priority_details, df, logo_path):
    """Generate PowerPoint presentation"""
    try:
        ppt_path = create_ppt_report(company_name, audit_date, kpis, exec_summary, priority_details, df, logo_path)
        return ppt_path
    except Exception as e:
        print(f"❌ PPT generation failed: {e}")
        return None

if __name__=="__main__":
    main()